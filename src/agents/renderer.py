"""
powerpoint rendering using python-pptx
"""

import re
import qrcode
from pathlib import Path
from typing import Dict, Any, Optional
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image

from src.state.poster_state import PosterState
from utils.src.logging_utils import log_agent_info, log_agent_success, log_agent_error
from src.config.poster_config import load_config


class Renderer:
    """powerpoint rendering with styling support"""
    
    def __init__(self):
        self.name = "renderer"
        self.styling_interfaces = None
        
        # load configuration
        self.config = load_config()
        self.layout_constants = self.config["layout_constants"]
        self.powerpoint_config = self.config["powerpoint"]
        self.indentation_config = self.config["indentation"]
        self.typography_config = self.config["typography"]

    def __call__(self, state: PosterState) -> PosterState:
        log_agent_info(self.name, "Starting Rendering Process")
        
        try:
            self.styling_interfaces = self._load_styling_interfaces(state)
            output_path = Path(state["output_dir"]) / f"{state['poster_name']}.pptx"
            self._render_presentation(state, output_path)
            
            # convert to png if possible
            png_path = self._convert_to_png(output_path)
            
            log_agent_success(self.name, f"rendered poster: {output_path}")
            if png_path:
                log_agent_success(self.name, f"generated preview: {png_path}")
                
        except Exception as e:
            log_agent_error(self.name, f"rendering failed: {e}")
            state["errors"].append(f"{self.name}: {e}")
            
        return state

    def _load_styling_interfaces(self, state: PosterState) -> Dict[str, Any]:
        """load styling interfaces from font agent output file"""
        styling_path = Path(state["output_dir"]) / "content" / "styling_interfaces.json"
        if styling_path.exists():
            with open(styling_path, 'r', encoding='utf-8') as f:
                interfaces = json.load(f)
            interfaces["line_spacing"] = 1.0
            return interfaces
        else:
            # fallback to defaults with 1.0 line spacing
            return {
                "bullet_point_marker": "•",
                "bold_start_tag": "**",
                "bold_end_tag": "**",
                "italic_start_tag": "*",
                "italic_end_tag": "*",
                "color_start_tag": "<color:",
                "color_end_tag": "</color>",
                "line_spacing": 1.0,
                "paragraph_spacing": 0.1
            }

    def _render_presentation(self, state: PosterState, output_path: Path):
        """render complete presentation"""
        prs = Presentation()
        prs.slide_width = Inches(state["poster_width"])
        prs.slide_height = Inches(state["poster_height"])
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # TODO: generate QR code if needed
        qr_code_path = None
        if state.get("url"):
            qr_code_path = self._generate_qr_code(state["url"], state["output_dir"])

        # use styled_layout if available, fallback to design_layout
        layout_data = state.get("styled_layout", state.get("design_layout", []))
        if not layout_data:
            raise ValueError("no styled_layout or design_layout found")
        
        # sort elements by priority to ensure proper rendering order
        sorted_elements = sorted(layout_data, key=lambda x: x.get("priority", 0.5))
        
        for element in sorted_elements:
            self._render_element(slide, element, state, qr_code_path)
        
        prs.save(output_path)

    def _render_element(self, slide, element: Dict, state: PosterState, qr_code_path: Optional[str]):
        """render individual element based on type"""
        element_type = element.get("type")
        
        # handle QR code elements
        if element_type == "qr_code" and qr_code_path:
            self._render_qr_code(slide, element, qr_code_path)
            return
        
        # get appropriate renderer
        renderer_map = {
            "title": self._render_title,
            "section_title": self._render_section_title,
            "title_accent_block": self._render_title_accent_block,
            "title_accent_line": self._render_title_accent_line,
            "conf_logo": self._render_conf_logo,
            "aff_logo": self._render_aff_logo,
            "section_container": self._render_section_container,
            "text": self._render_text,
            "visual": self._render_visual,
            "mixed": self._render_mixed,
        }
        
        renderer = renderer_map.get(element_type)
        if renderer:
            renderer(slide, element, state)
        else:
            log_agent_error(self.name, f"unknown element type: {element_type}")

    def _render_title(self, slide, element: Dict, state: PosterState):
        """render poster title with authors"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        log_agent_info(self.name, f"rendering title at ({x.inches:.1f}, {y.inches:.1f})")
        
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Make sure title fits in the fixed-height textbox
        tf.word_wrap = True
        
        content = element.get("content", "Title\nAuthors")
        lines = content.split("\n")
        
        # separate title and authors
        title_lines = lines[:-1] if len(lines) > 1 else lines
        authors_text = lines[-1] if len(lines) > 1 else ""
        
        # add title lines
        for i, title_line in enumerate(title_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = title_line.strip()
            p.font.name = element.get("font_family", "Helvetica Neue")
            title_font_size = self.styling_interfaces.get("font_sizes", {}).get("title", 100)
            p.font.size = Pt(element.get("font_size", title_font_size))
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # black for readability
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = self.typography_config["line_spacing"]
        
        # add authors
        if authors_text:
            p_authors = tf.add_paragraph()
            p_authors.text = authors_text.strip()
            p_authors.font.name = "Arial"
            authors_font_size = self.styling_interfaces.get("font_sizes", {}).get("authors", 72)
            p_authors.font.size = Pt(element.get("author_font_size", authors_font_size))
            p_authors.font.color.rgb = RGBColor(60, 60, 60)  # dark gray
            p_authors.alignment = PP_ALIGN.LEFT
            p_authors.line_spacing = self.typography_config["line_spacing"] + 0.1  # slightly looser for authors

    def _render_section_title(self, slide, element: Dict, state: PosterState):
        """render section title with enhanced styling"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        section_title = element.get("section_title", "").strip()
        if not section_title:
            return
        
        log_agent_info(self.name, f"rendering section title: '{section_title}'")
        
        # create textbox for section title
        textbox = slide.shapes.add_textbox(x, y, w, h)
        tf = textbox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = False
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        
        # use existing first paragraph to avoid extra newline
        if len(tf.paragraphs) > 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = section_title
        p.font.name = element.get("font_family", "Helvetica Neue")
        section_title_font_size = self.styling_interfaces.get("font_sizes", {}).get("section_title", 48)
        p.font.size = Pt(element.get("font_size", section_title_font_size))
        p.font.bold = element.get("font_weight", "bold") == "bold"
        
        # apply color styling
        font_color = element.get("font_color", "#000000")
        p.font.color.rgb = self._parse_color(font_color)
        
        # apply alignment based on design template
        alignment = element.get("alignment", "left").lower()
        if alignment == "center":
            p.alignment = PP_ALIGN.CENTER
        elif alignment == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

    def _render_title_accent_block(self, slide, element: Dict, state: PosterState):
        """render color block accent for section titles"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        # use 'color' field from layout agent
        fill_color = element.get("color", element.get("fill_color", "#1E3A8A"))
        
        log_agent_info(self.name, f"rendering title accent block: {fill_color} at ({x.inches:.2f}, {y.inches:.2f})")
        
        # create rectangle shape
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = self._parse_color(fill_color)
        rect.line.fill.background()  # no border

    def _render_title_accent_line(self, slide, element: Dict, state: PosterState):
        """render underline accent for section titles"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        # use 'color' field from layout agent
        fill_color = element.get("color", element.get("fill_color", "#E8E8E8"))
        
        log_agent_info(self.name, f"rendering title accent line: {fill_color} at ({x.inches:.2f}, {y.inches:.2f})")
        
        # create thin rectangle for line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        line.fill.solid()
        line.fill.fore_color.rgb = self._parse_color(fill_color)
        line.line.fill.background()  # no border

    def _render_section_container(self, slide, element: Dict, state: PosterState):
        """render section container with optional debug border and mono_light background for critical sections"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        is_debug = element.get("debug_border", False)
        importance_level = element.get("importance_level", 2)
        
        # create base rectangle
        container = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        
        # apply background fill based on importance level
        if importance_level == 1:
            # critical section gets mono_light background color
            color_scheme = state.get("color_scheme", {})
            mono_light = color_scheme.get("mono_light", "#e6eaef")
            container.fill.solid()
            container.fill.fore_color.rgb = self._parse_color(mono_light)
            log_agent_info(self.name, f"applied mono_light background ({mono_light}) to critical section")
        else:
            # non-critical sections remain transparent
            container.fill.background()
        
        # apply border based on debug mode
        if is_debug:
            # prominent debug border
            container.line.color.rgb = RGBColor(255, 0, 0)  # red border
            container.line.width = Pt(2)
            log_agent_info(self.name, f"added debug section border")
        else:
            container.line.fill.background()

    def _render_text(self, slide, element: Dict, state: PosterState):
        """render text elements with enhanced formatting"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        content = element.get("content", "").strip()
        if not content:
            return
        
        log_agent_info(self.name, f"rendering text element: {element.get('id', 'unknown')}")
        
        # add text with margins
        margin = self.layout_constants["text_margin_renderer"]  # reduced margin for better space utilization
        self._add_enhanced_text(
            slide, content, 
            x + Inches(margin), y, 
            w - Inches(2 * margin), h,
            element
        )

    def _render_visual(self, slide, element: Dict, state: PosterState):
        """render visual elements with proper aspect ratio and scaling"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        visual_id = element.get("visual_id")
        scale_factor = element.get("scale_factor", 1.0)  # default to no scaling
        
        if visual_id:
            # layout agent already calculated padding, use exact positioning
            self._add_visual_with_aspect_ratio(
                slide, visual_id, 
                x, y, w, h,
                state, scale_factor
            )

    def _render_mixed(self, slide, element: Dict, state: PosterState):
        """render mixed elements (text and visual)"""
        # for now, treat as text element
        self._render_text(slide, element, state)

    def _render_conf_logo(self, slide, element: Dict, state: PosterState):
        """render conference logo"""
        logo_path = state.get("logo_path")
        if logo_path and Path(logo_path).exists():
            self._render_logo_with_aspect_ratio(slide, element, logo_path)

    def _render_aff_logo(self, slide, element: Dict, state: PosterState):
        """render affiliation logo"""
        aff_logo_path = state.get("aff_logo_path")
        if aff_logo_path and Path(aff_logo_path).exists():
            self._render_logo_with_aspect_ratio(slide, element, aff_logo_path)

    def _add_enhanced_text(self, slide, text: str, left, top, width, height, element: Dict):
        """add text with enhanced formatting support"""
        if not text.strip():
            return
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.clear()
        
        # enforce height constraints to prevent text overflow beyond textbox bounds
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        # add small margins to ensure text stays within bounds
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        # get font properties from element
        font_family = element.get("font_family", "Arial")
        font_size = element.get("font_size", 40)
        font_color = element.get("font_color", "#000000")
        line_spacing = element.get("line_spacing", self.styling_interfaces["line_spacing"])
        
        self._format_enhanced_text(tf, text, font_family, font_size, font_color, line_spacing)
        
        # debug info for formatting
        total_runs = sum(len(p.runs) for p in tf.paragraphs)
        log_agent_info(self.name, f"created {len(tf.paragraphs)} paragraphs with {total_runs} formatted runs")

    def _format_enhanced_text(self, text_frame, text: str, font_family: str, font_size: int, font_color: str, line_spacing: float):
        """format text with enhanced bullet point and bold support using 1.0 line spacing"""
        text_frame.clear()
        
        body_text_font_size = self.styling_interfaces.get("font_sizes", {}).get("body_text", 40)
        effective_font_size = font_size if font_size != 40 else body_text_font_size
        base_font_size = Pt(max(effective_font_size, 36))  # minimum 36pt
        base_color = self._parse_color(font_color)
        
        # split by single newlines only (treat as simple line breaks)
        lines = text.split('\n')
        
        for line_idx, line in enumerate(lines):
            original_line = line  # keep original line for indentation detection
            line = line.strip()
            if not line:
                continue
            
            # create paragraph for each line
            if line_idx == 0 and len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # handle indentation by checking if line starts with ◦ (sub-bullet)
            if line.strip().startswith(self.indentation_config["secondary_bullet_char"]):  # secondary bullet character
                # set paragraph level for indentation
                p.level = self.indentation_config["secondary_level"]
            else:
                p.level = self.indentation_config["primary_level"]
            
            # add formatted text content (don't clear p.text)
            self._add_formatted_runs(p, line, font_family, base_font_size, base_color)
            
            # set paragraph properties - force 1.0 line spacing
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = self.typography_config["line_spacing"]  # fixed 1.0 line spacing

    def _add_formatted_runs(self, paragraph, text: str, font_family: str, 
                          base_font_size, base_color):
        """add text with all formatting as separate runs - following pptx best practices"""
        self._parse_and_add_runs(paragraph, text, font_family, base_font_size, base_color)

    def _parse_and_add_runs(self, paragraph, text: str, font_family: str, 
                          base_font_size, base_color):
        """parse text and create separate runs for each format type"""
        # tokenize the text into formatting segments
        segments = self._tokenize_formatting(text)
        
        # create runs for each segment
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment['text']
            run.font.name = font_family
            run.font.size = base_font_size
            
            # apply formatting based on segment type
            if segment['color']:
                run.font.color.rgb = self._parse_color(segment['color'])
            else:
                run.font.color.rgb = base_color
            
            if segment['bold']:
                run.font.bold = True
            
            if segment['italic']:
                run.font.italic = True
    
    def _tokenize_formatting(self, text: str) -> list:
        """tokenize text into formatting segments with precise position tracking"""
        segments = []
        i = 0
        
        while i < len(text):
            # check for color markup: <color:#RRGGBB>text</color>
            color_match = re.match(r'<color:(#[0-9A-Fa-f]{6})>', text[i:])
            if color_match:
                color_hex = color_match.group(1)
                opening_tag_end = i + color_match.end()
                
                # find closing </color> tag using absolute position
                closing_tag_pattern = r'</color>'
                color_content_start = opening_tag_end
                closing_match = re.search(closing_tag_pattern, text[color_content_start:])
                
                if closing_match:
                    # calculate absolute positions
                    color_content_end = color_content_start + closing_match.start()
                    closing_tag_end = color_content_start + closing_match.end()
                    
                    # extract content between color tags
                    colored_text = text[color_content_start:color_content_end]
                    
                    # process colored text with automatic bold
                    if colored_text.strip():  # only process non-empty content
                        segments.append({
                            'text': colored_text,
                            'bold': True,  # all colored text is bold
                            'italic': False,
                            'color': color_hex
                        })
                    
                    # move past the entire color block
                    i = closing_tag_end
                    continue
                else:
                    # malformed color tag, treat as regular text
                    segments.append({
                        'text': text[i],
                        'bold': False,
                        'italic': False,
                        'color': None
                    })
                    i += 1
                    continue
            
            # check for bold: **text**
            bold_match = re.match(r'\*\*(.*?)\*\*', text[i:])
            if bold_match:
                bold_text = bold_match.group(1)
                segments.append({
                    'text': bold_text,
                    'bold': True,
                    'italic': False,
                    'color': None
                })
                i += bold_match.end()
                continue
            
            # check for italic: *text*
            italic_match = re.match(r'\*(.*?)\*', text[i:])
            if italic_match:
                italic_text = italic_match.group(1)
                segments.append({
                    'text': italic_text,
                    'bold': False,
                    'italic': True,
                    'color': None
                })
                i += italic_match.end()
                continue
            
            # regular text - find next formatting marker
            next_format = re.search(r'(\*\*|\*|<color:)', text[i:])
            if next_format:
                regular_text = text[i:i + next_format.start()]
            else:
                regular_text = text[i:]
            
            if regular_text:
                segments.append({
                    'text': regular_text,
                    'bold': False,
                    'italic': False,
                    'color': None
                })
            
            if next_format:
                i += next_format.start()
            else:
                break
        
        return segments
    
    def _parse_bold_italic(self, text: str, color: str) -> list:
        """simplified bold/italic parser - only used for nested formatting"""
        segments = []
        i = 0
        
        while i < len(text):
            # check for bold
            bold_match = re.match(r'\*\*(.*?)\*\*', text[i:])
            if bold_match:
                bold_text = bold_match.group(1)
                segments.append({
                    'text': bold_text,
                    'bold': True,
                    'italic': False,
                    'color': color
                })
                i += bold_match.end()
                continue
            
            # check for italic
            italic_match = re.match(r'\*(.*?)\*', text[i:])
            if italic_match:
                italic_text = italic_match.group(1)
                segments.append({
                    'text': italic_text,
                    'bold': bool(color),  # force bold if color is present
                    'italic': True,
                    'color': color
                })
                i += italic_match.end()
                continue
            
            # regular text
            next_format = re.search(r'(\*\*|\*)', text[i:])
            if next_format:
                regular_text = text[i:i + next_format.start()]
            else:
                regular_text = text[i:]
            
            if regular_text:
                segments.append({
                    'text': regular_text,
                    'bold': bool(color),  # force bold if color is present
                    'italic': False,
                    'color': color
                })
            
            if next_format:
                i += next_format.start()
            else:
                break
        
        return segments

    def _add_visual_with_aspect_ratio(self, slide, visual_id: str, left, top, width, height, state, scale_factor: float = 1.0):
        """add visual with proper aspect ratio preservation and optional scaling"""
        visual_path = self._get_visual_path(visual_id, state)
        
        if visual_path and Path(visual_path).exists():
            try:
                # calculate proper size maintaining aspect ratio
                with Image.open(visual_path) as img:
                    orig_width, orig_height = img.size
                    aspect_ratio = orig_width / orig_height
                
                # get allocated space from layout
                available_width = width.inches if hasattr(width, 'inches') else float(width)
                available_height = height.inches if hasattr(height, 'inches') else float(height)
                
                # always use exact dimensions and positioning from JSON
                final_width = Inches(available_width)
                final_height = Inches(available_height)
                centered_left = left
                centered_top = top
                
                slide.shapes.add_picture(visual_path, centered_left, centered_top, width=final_width, height=final_height)
                
                if scale_factor < 1.0:
                    log_agent_info(self.name, f"visual {visual_id} uses layout-calculated dimensions (scale_factor={scale_factor:.1f} already applied)")
                                       
            except Exception as e:
                log_agent_error(self.name, f"failed to add visual {visual_id}: {e}")

    def _render_logo_with_aspect_ratio(self, slide, element: Dict, image_path: str):
        """render logo with proper aspect ratio preservation"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        
        try:
            # calculate dimensions while preserving aspect ratio
            with Image.open(image_path) as img:
                orig_width, orig_height = img.size
                aspect_ratio = orig_width / orig_height
            
            available_width = w.inches if hasattr(w, 'inches') else float(w)
            available_height = h.inches if hasattr(h, 'inches') else float(h)
            
            # fit image within available space
            if available_width / aspect_ratio <= available_height:
                final_width = Inches(available_width)
                final_height = Inches(available_width / aspect_ratio)
            else:
                final_height = Inches(available_height)
                final_width = Inches(available_height * aspect_ratio)
            
            # center the image
            centered_left = x + (w - final_width) / 2
            centered_top = y + (h - final_height) / 2
            
            slide.shapes.add_picture(image_path, centered_left, centered_top, 
                                   width=final_width, height=final_height)
                                   
        except Exception as e:
            log_agent_error(self.name, f"failed to render logo: {e}")

    def _get_visual_path(self, visual_id: str, state: PosterState) -> Optional[str]:
        """get path to visual asset"""
        images = state.get("images", {})
        tables = state.get("tables", {})
        vid = (visual_id or "").split('_')[-1]
        
        if visual_id.startswith("figure"):
            return images.get(vid, {}).get("path")
        if visual_id.startswith("table"):
            return tables.get(vid, {}).get("path")
        
        return None

    def _parse_color(self, color_str: str) -> RGBColor:
        """parse color string to RGBColor"""
        hex_color = color_str.lstrip('#')
        r, g, b = (int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)

    def _generate_qr_code(self, url: str, output_dir: str) -> str:
        """generate QR code for URL"""
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=2,
        )
        qr.add_data(url)
        qr.make(fit=True)
        
        img = qr.make_image(fill_color="black", back_color="white")
        qr_path = Path(output_dir) / "qr_code.png"
        img.save(qr_path)
        
        return str(qr_path)

    def _render_qr_code(self, slide, element: Dict, qr_code_path: str):
        """render QR code element"""
        x, y, w, h = (Inches(element[k]) for k in ["x", "y", "width", "height"])
        slide.shapes.add_picture(qr_code_path, x, y, w, h)

    def _convert_to_png(self, pptx_path: Path) -> Optional[str]:
        """convert PPTX to PNG using LibreOffice"""
        try:
            import subprocess
            output_dir = pptx_path.parent
            
            import platform
            system = platform.system().lower()
            
            if system == "windows":
                libreoffice_paths = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    r"C:\Users\%USERNAME%\AppData\Local\Programs\LibreOffice\program\soffice.exe",
                    "soffice.exe",
                    "libreoffice.exe"
                ]
            elif system == "linux":
                libreoffice_paths = [
                    "/usr/bin/libreoffice",
                    "/usr/local/bin/libreoffice",
                    "/snap/bin/libreoffice",
                    "/usr/bin/soffice",
                    "libreoffice",
                    "soffice"
                ]
            elif system == "darwin":  # macOS
                libreoffice_paths = [
                    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                    "/usr/local/bin/libreoffice",
                    "libreoffice",
                    "soffice"
                ]
            else:
                libreoffice_paths = [
                    "libreoffice",
                    "soffice"
                ]
            
            for lo_path in libreoffice_paths:
                try:
                    cmd = [
                        lo_path, "--headless", "--convert-to", "png",
                        "--outdir", str(output_dir), str(pptx_path)
                    ]
                    
                    result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
                    
                    if result.returncode == 0:
                        png_name = pptx_path.stem + ".png"
                        png_path = output_dir / png_name
                        if png_path.exists():
                            return str(png_path)
                            
                except (subprocess.SubprocessError, FileNotFoundError):
                    continue
            
            log_agent_error(self.name, "LibreOffice not found - install for PNG conversion")
            
        except Exception as e:
            log_agent_error(self.name, f"PNG conversion failed: {e}")
            
        return None


def renderer_node(state: PosterState) -> Dict[str, Any]:
    result = Renderer()(state)
    return {
        **state,
        "tokens": result["tokens"],
        "current_agent": result["current_agent"],
        "errors": result["errors"]
    }