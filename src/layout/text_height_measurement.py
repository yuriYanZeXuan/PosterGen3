"""
text height measurement using binary search overflow detection
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from typing import Dict, Any, List
from pathlib import Path
from src.config.poster_config import load_config

def get_font_file_path(font_name: str) -> str:
    font_mapping = {
        "Arial": "fonts/Arial.ttf",
        "Helvetica Neue": "fonts/HelveticaNeue.ttf",
    }
    
    font_file = font_mapping.get(font_name, "fonts/Arial.ttf")
    project_root = Path(__file__).parent.parent.parent
    font_path = project_root / font_file
    
    return str(font_path)

def measure_text_height(text_content: str, width_inches: float, font_name: str = "Arial", 
                       font_size: int = 44, line_spacing: float = 1.0, precision: float = 0.001) -> Dict[str, Any]:
    """find minimum height for text to fit without font size reduction"""
    
    config = load_config()
    
    prs = Presentation()
    config = load_config()
    slide_layout_index = config["powerpoint"]["slide_layout_blank"]
    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
    
    min_height = config["text_measurement"]["min_height"]
    max_height = config["text_measurement"]["max_height"]
    tolerance = precision
    
    while (max_height - min_height) > tolerance:
        test_height = (min_height + max_height) / 2
        
        textbox = slide.shapes.add_textbox(
            left=Inches(config["powerpoint"]["text_frame_positioning"]["default_left"]),
            top=Inches(config["powerpoint"]["text_frame_positioning"]["default_top"]),
            width=Inches(width_inches),
            height=Inches(test_height)
        )
        
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        # use same margins as layout agent for consistent measurement
        text_frame.margin_left = Inches(config["text_measurement"]["margins"]["left"])
        text_frame.margin_right = Inches(config["text_measurement"]["margins"]["right"])
        text_frame.margin_top = Inches(config["text_measurement"]["margins"]["top"])
        text_frame.margin_bottom = Inches(config["text_measurement"]["margins"]["bottom"])
        
        # process text exactly like renderer: split by single newlines
        lines = text_content.split('\n')
        
        for line_idx, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            
            # create paragraph for each line (matching renderer behavior)
            if line_idx == 0 and len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = line
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = line_spacing
            
            # apply font to each paragraph
            if p.runs:
                run = p.runs[0]
                run.font.name = font_name
                run.font.size = Pt(font_size)
        
        original_size = font_size
        font_reduced = False
        
        try:
            # Use direct font file to bypass cross-platform discovery bug
            font_file_path = get_font_file_path(font_name)
            text_frame.fit_text(font_file=font_file_path, max_size=font_size)
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size and run.font.size.pt < (original_size - 0.5):
                        font_reduced = True
                        break
                if font_reduced:
                    break
        except Exception as e:
            print(f"fit_text error: {e}")
            font_reduced = True
        
        if font_reduced:
            min_height = test_height
        else:
            max_height = test_height
        
        # cleanup textbox
        sp = textbox._element
        sp.getparent().remove(sp)
    
    # calculate newline offset to compensate for pptx rendering discrepancy
    newline_count = text_content.count('\n')
    newline_offset = newline_count * (font_size / 72) * config["text_measurement"]["newline_offset_ratio"]
    final_height = max_height + newline_offset
    
    return {
        "optimal_height": final_height,
        "text_content": text_content,
        "width_inches": width_inches,
        "font_name": font_name,
        "font_size": font_size,
        "line_spacing": line_spacing,
        "precision": precision,
        "newline_count": newline_count,
        "newline_offset": newline_offset
    }